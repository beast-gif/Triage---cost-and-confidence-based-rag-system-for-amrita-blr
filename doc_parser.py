"""
doc_parser.py — turn any uploaded document into chunks.

ONE GENERIC PATH, NOT A PARSER PER FORMAT
------------------------------------------
A calendar PDF, a lecture deck and a policy document share no structure. So
this does not try to understand documents — it extracts BLOCKS of text with
their page or slide number, and everything downstream is format-agnostic:

    file -> extract_blocks() -> chunk_blocks() -> chunks

extract_blocks returns a flat list of {text, page, kind} whatever the source.
Adding a format means adding one extractor; nothing else changes.

WHY TABLES GET A PROSE PREFIX
-----------------------------
extractor.py learned this the hard way on the web corpus: raw pipe-separated
table text embeds badly against natural-language questions. A fee table only
started matching "what is the fee for X" once a framing sentence was prepended.
Same applies here — every table block gets "From {title}, page {n}:" in front.

WHY pdfplumber AND NOT pypdf
----------------------------
Measured on the academic calendar. pypdf reads by text-object order, so events
in merged cells were dumped at the bottom of the page, detached from their
dates:

    pypdf:       "Commencement of classes for UG-S3..."   (no date anywhere)
    pdfplumber:  ['08-Jul', 'Wed', '1', '', 'Commencement of classes...']

The pypdf version would retrieve on "when do classes commence" and contain no
answer — the same lexically-right, semantically-empty failure the confidence
score cannot detect.

Usage:
    python doc_parser.py calender.pdf
    python doc_parser.py deck.pptx --verbose
"""

import os
import re

# Chunk sizing. Blocks below MIN_CHARS get merged with their neighbour — a
# lone "01-Jul Wed" row carries no meaning on its own and retrieves badly.
MAX_CHARS = 1200
MIN_CHARS = 120

SUPPORTED = {".pdf", ".pptx", ".docx", ".txt", ".md"}


def _clean(text):
    """Collapse the whitespace PDF and PPTX extraction leaves behind."""
    if not text:
        return ""
    text = text.replace("\xa0", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _table_to_text(rows):
    """
    Pipe-separate a table, dropping empty cells and empty rows.

    pdfplumber returns None for cells that were never populated and '' for
    cells that exist but are blank; both become nothing here.
    """
    lines = []
    for row in rows or []:
        cells = [_clean(str(c)).replace("\n", " ") for c in row if c]
        cells = [c for c in cells if c]
        if cells:
            lines.append(" | ".join(cells))
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# per-format extractors — each returns [{text, page, kind}]
# ---------------------------------------------------------------------------
def _extract_pdf(path):
    import pdfplumber

    blocks = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            tables = page.extract_tables() or []
            table_text = ""
            for rows in tables:
                text = _table_to_text(rows)
                if text:
                    blocks.append({"text": text, "page": i, "kind": "table"})
                    table_text += text

            prose = _clean(page.extract_text() or "")
            if not prose:
                continue

            # On a page that is mostly table, extract_text() returns the SAME
            # content without the cell separators — '01-Jul Wed' instead of
            # '01-Jul | Wed'. Storing both duplicates every row and gives the
            # worse version equal footing in retrieval.
            #
            # Measured on the academic calendar: taking both produced 12 blocks
            # for 6 pages, every page indexed twice.
            #
            # So prose is only kept when it carries meaningfully MORE than the
            # tables did — headings, footnotes, callouts sitting outside the
            # grid. The ratio is deliberately generous: dropping a real
            # footnote is worse than keeping a little overlap.
            if table_text and len(prose) < len(table_text) * 1.25:
                continue

            blocks.append({"text": prose, "page": i, "kind": "text"})
    return blocks


def _extract_pptx(path):
    from pptx import Presentation

    blocks = []
    for i, slide in enumerate(Presentation(path).slides, start=1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = _clean(shape.text_frame.text)
                if t:
                    parts.append(t)
            if getattr(shape, "has_table", False):
                rows = [[c.text for c in r.cells] for r in shape.table.rows]
                t = _table_to_text(rows)
                if t:
                    parts.append(t)

        # Speaker notes carry the explanation the slide only gestures at, so
        # they are worth indexing.
        if slide.has_notes_slide:
            notes = _clean(slide.notes_slide.notes_text_frame.text)
            if notes:
                parts.append(f"Speaker notes: {notes}")

        if parts:
            blocks.append({"text": "\n".join(parts), "page": i, "kind": "slide"})
    return blocks


def _extract_docx(path):
    import docx

    document = docx.Document(path)
    blocks = []

    paragraphs = [_clean(p.text) for p in document.paragraphs]
    body = "\n".join(p for p in paragraphs if p)
    if body:
        blocks.append({"text": body, "page": 1, "kind": "text"})

    for table in document.tables:
        rows = [[c.text for c in r.cells] for r in table.rows]
        t = _table_to_text(rows)
        if t:
            blocks.append({"text": t, "page": 1, "kind": "table"})
    return blocks


def _extract_plain(path):
    with open(path, encoding="utf-8", errors="replace") as f:
        raw = f.read()
    parts = [_clean(p) for p in re.split(r"\n\s*\n", raw)]
    return [{"text": p, "page": 1, "kind": "text"} for p in parts if p]


EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".pptx": _extract_pptx,
    ".docx": _extract_docx,
    ".txt": _extract_plain,
    ".md": _extract_plain,
}


def extract_blocks(path):
    ext = os.path.splitext(path)[1].lower()
    if ext not in EXTRACTORS:
        raise ValueError(
            f"unsupported file type '{ext}' — supported: {', '.join(sorted(SUPPORTED))}"
        )
    return EXTRACTORS[ext](path)


# ---------------------------------------------------------------------------
# chunking
# ---------------------------------------------------------------------------
def _split_long(text, limit=MAX_CHARS):
    """Split on paragraph breaks, then lines, never mid-word."""
    if len(text) <= limit:
        return [text]

    pieces, current = [], ""
    for para in re.split(r"\n\s*\n", text):
        if len(current) + len(para) + 2 <= limit:
            current = f"{current}\n\n{para}" if current else para
            continue
        if current:
            pieces.append(current)
        if len(para) <= limit:
            current = para
            continue
        # A single oversized paragraph — fall back to line boundaries.
        current = ""
        for line in para.split("\n"):
            if len(current) + len(line) + 1 <= limit:
                current = f"{current}\n{line}" if current else line
            else:
                if current:
                    pieces.append(current)
                current = line[:limit]
    if current:
        pieces.append(current)
    return pieces


def chunk_blocks(blocks, doc_id, title):
    """
    Blocks -> chunks, with a prose header on every chunk.

    The header is not decoration. A bare table row embeds poorly against a
    natural-language question; naming the document and page gives the embedding
    model something to match on, and gives the answer generator a citation.
    """
    chunks, buffer, buffer_page = [], "", None

    def flush():
        nonlocal buffer, buffer_page
        if not buffer.strip():
            buffer = ""
            return
        pieces = _split_long(buffer)

        # The first line of a block usually carries its identity — a month
        # header, a table's column row, a slide title. Splitting leaves every
        # piece after the first without it.
        #
        # Measured on the academic calendar: part 2 of the July page held
        # '29-Jul | Wed | H | Guru Poornima (Holiday)' with no year and no
        # indication it was a calendar at all. Retrievable, but the generator
        # would have to guess the year.
        #
        # So the first line rides along on every continuation. A few repeated
        # tokens buys chunks that stand alone.
        context_line = pieces[0].split("\n", 1)[0].strip() if len(pieces) > 1 else ""

        for n, piece in enumerate(pieces):
            index = len(chunks)
            # The header goes on EVERY piece, not just the first. A buffer that
            # splits into three chunks used to leave two of them headerless —
            # no document name, no page number, nothing to cite.
            header = f"From {title}, page {buffer_page}:"
            if len(pieces) > 1:
                header += f" (part {n + 1} of {len(pieces)})"
            body = piece
            if n > 0 and context_line and context_line not in piece:
                body = f"{context_line}\n{piece}"
            chunks.append({
                "chunk_id": f"upload:{doc_id}:{index}",
                "doc_id": doc_id,
                "title": title,
                "page": buffer_page,
                "source_type": "upload",
                "content": f"{header}\n\n{body}",
            })
        buffer = ""

    for block in blocks:
        text = block["text"]
        if buffer_page is not None and block["page"] != buffer_page:
            flush()
        buffer_page = block["page"]

        # Merge blocks that are too small to stand alone, rather than emitting
        # a chunk that is mostly header.
        if buffer and (len(buffer) < MIN_CHARS or len(buffer) + len(text) <= MAX_CHARS):
            buffer = f"{buffer}\n\n{text}"
        else:
            flush()
            buffer = text

    flush()
    return chunks


def parse_document(path, doc_id=None, title=None):
    """file -> chunks. doc_id and title default to the filename."""
    name = os.path.basename(path)
    doc_id = doc_id or re.sub(r"[^a-z0-9]+", "-", os.path.splitext(name)[0].lower()).strip("-")
    title = title or os.path.splitext(name)[0]
    return chunk_blocks(extract_blocks(path), doc_id, title)


if __name__ == "__main__":
    import sys

    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    if not args:
        print("usage: python doc_parser.py <file> [--verbose]")
        sys.exit(1)

    path = args[0]
    verbose = "--verbose" in sys.argv

    blocks = extract_blocks(path)
    chunks = parse_document(path)

    kinds = {}
    for b in blocks:
        kinds[b["kind"]] = kinds.get(b["kind"], 0) + 1

    print(f"\n{path}")
    print("=" * 70)
    print(f"blocks : {len(blocks)}  {kinds}")
    print(f"chunks : {len(chunks)}")
    if chunks:
        sizes = [len(c["content"]) for c in chunks]
        print(f"chars  : min={min(sizes)}  mean={sum(sizes)//len(sizes)}  max={max(sizes)}")

    show = chunks if verbose else chunks[:3]
    for c in show:
        print(f"\n--- {c['chunk_id']}  (page {c['page']}, {len(c['content'])} chars)")
        body = c["content"] if verbose else c["content"][:400]
        print(body)